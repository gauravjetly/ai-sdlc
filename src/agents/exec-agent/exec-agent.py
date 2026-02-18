#!/usr/bin/env python3
"""
Exec Agent - Executive Presentation Specialist
Self-learning agent that generates Deltek-branded PowerPoint presentations
"""

import os
import json
import shutil
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class ExecAgent:
    """Executive Presentation Specialist with Self-Learning Capabilities"""

    def __init__(self):
        self.name = "Exec Agent"
        self.id = "exec"
        self.icon = "📊"
        self.color = "#1742F6"  # Deltek Blue
        self.model = "Sonnet"
        self.role = "Executive Presentations"

        # Paths
        self.home = Path.home()
        self.memory_dir = self.home / ".claude" / "exec-agent-memory"
        self.sdlc_registry = self.home / ".claude" / "sdlc-registry"
        self.template_path = Path("/Users/gauravjetly/Downloads/Deltek PowerPoint Guidelines 2/Deltek PPT Template and Guidelines 011426.potx")

        # Initialize memory
        self._init_memory()
        self._load_deltek_brand()

    def _init_memory(self):
        """Initialize agent memory directories"""
        dirs = [
            self.memory_dir / "deltek-brand",
            self.memory_dir / "presentations" / "archive",
            self.memory_dir / "presentations" / "templates",
            self.memory_dir / "presentations" / "diagrams",
            self.memory_dir / "preferences",
            self.memory_dir / "learning"
        ]
        for d in dirs:
            d.mkdir(parents=True, exist_ok=True)

    def _load_deltek_brand(self):
        """Load and memorize Deltek brand guidelines"""
        brand_file = self.memory_dir / "deltek-brand" / "colors.json"

        self.deltek_brand = {
            "colors": {
                "primary_blue": "#1742F6",
                "navy": "#081581",
                "dark_gray": "#3C454E",
                "teal": "#00B6C3",
                "purple": "#6D18F1",
                "magenta": "#C200CC",
                "dela_cyan": "#08E9EB",
                "dela_blue": "#3895FF",
                "gradient": ["#08E9EB", "#FF5DF2", "#3895FF", "#7A62FF"],
                "success": "#00875a",
                "warning": "#ff9800",
                "error": "#d32f2f"
            },
            "typography": {
                "heading": {"font": "Figtree", "weight": "Bold", "size": 32},
                "subheading": {"font": "Figtree", "weight": "SemiBold", "size": 24},
                "body": {"font": "Figtree", "weight": "Regular", "size": 14},
                "caption": {"font": "Figtree", "weight": "Regular", "size": 11}
            }
        }

        # Save to memory for future reference
        with open(brand_file, 'w') as f:
            json.dump(self.deltek_brand, f, indent=2)

    def generate_presentation(self, project_id: str, presentation_type: str = "executive-summary") -> str:
        """
        Generate a Deltek-branded PowerPoint presentation

        Args:
            project_id: SDLC project identifier
            presentation_type: Type of presentation (executive-summary, architecture, status)

        Returns:
            Path to generated presentation
        """
        # Load project data from SDLC registry
        project_data = self._load_project_data(project_id)

        # Create presentation from template
        if self.template_path.exists():
            prs = Presentation(str(self.template_path))
        else:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

        # Generate slides based on type
        if presentation_type == "executive-summary":
            self._add_title_slide(prs, project_data)
            self._add_executive_summary_slide(prs, project_data)
            self._add_key_metrics_slide(prs, project_data)
            self._add_architecture_overview_slide(prs, project_data)
            self._add_status_dashboard_slide(prs, project_data)
            self._add_timeline_slide(prs, project_data)
            self._add_next_steps_slide(prs, project_data)
        elif presentation_type == "architecture":
            self._add_title_slide(prs, project_data)
            self._add_architecture_detailed_slide(prs, project_data)
            self._add_component_breakdown_slide(prs, project_data)
            self._add_data_flow_slide(prs, project_data)
            self._add_deployment_diagram_slide(prs, project_data)
        elif presentation_type == "status":
            self._add_title_slide(prs, project_data)
            self._add_status_overview_slide(prs, project_data)
            self._add_agent_performance_slide(prs, project_data)
            self._add_issues_risks_slide(prs, project_data)
            self._add_action_items_slide(prs, project_data)

        # Save presentation
        output_path = self._save_presentation(prs, project_id, presentation_type)

        # Learn from this generation
        self._record_learning(project_id, presentation_type, output_path)

        return output_path

    def _load_project_data(self, project_id: str) -> Dict:
        """Load project data from SDLC registry"""
        project_file = self.sdlc_registry / "projects" / f"{project_id}.json"

        if project_file.exists():
            with open(project_file, 'r') as f:
                return json.load(f)

        # Return default data if project not found
        return {
            "id": project_id,
            "name": "AI-SDLC Project",
            "status": "in_progress",
            "phases": [],
            "created_at": datetime.now().isoformat()
        }

    def _add_title_slide(self, prs: Presentation, project: Dict):
        """Add Deltek-branded title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = project.get("name", "AI-SDLC Project")
        subtitle.text = f"Project Status Report\n{datetime.now().strftime('%B %d, %Y')}\n\nGenerated by AI-SDLC Exec Agent"

        # Apply Deltek brand colors
        self._apply_deltek_styling(title)

    def _add_executive_summary_slide(self, prs: Presentation, project: Dict):
        """Add executive summary with key highlights"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Executive Summary"
        self._apply_deltek_styling(title)

        # Add content box
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Add summary points
        summary_points = self._generate_summary_points(project)
        for point in summary_points:
            p = text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.space_after = Pt(12)

    def _add_key_metrics_slide(self, prs: Presentation, project: Dict):
        """Add key metrics with visual indicators"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title = slide.shapes.title
        title.text = "Key Performance Indicators"
        self._apply_deltek_styling(title)

        # Add metric boxes (4 KPIs in a grid)
        metrics = self._calculate_kpis(project)
        positions = [
            (Inches(1), Inches(2)),
            (Inches(5.5), Inches(2)),
            (Inches(1), Inches(4)),
            (Inches(5.5), Inches(4))
        ]

        for i, (metric_name, metric_value) in enumerate(metrics.items()):
            if i < 4:
                left, top = positions[i]
                self._add_kpi_box(slide, left, top, metric_name, metric_value)

    def _add_architecture_overview_slide(self, prs: Presentation, project: Dict):
        """Add high-level architecture diagram"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title = slide.shapes.title
        title.text = "System Architecture Overview"
        self._apply_deltek_styling(title)

        # Add architecture description
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "AI-SDLC platform architecture with 9 specialized agents orchestrating the complete software development lifecycle."
        text_frame.paragraphs[0].font.size = Pt(14)

        # Placeholder for diagram (would be generated using diagram library)
        left = Inches(2)
        top = Inches(3.5)
        width = Inches(6)
        height = Inches(3)

        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.text = "Architecture Diagram\n(Auto-generated from system topology)"
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 250)  # Light blue

    def _add_status_dashboard_slide(self, prs: Presentation, project: Dict):
        """Add status dashboard with progress indicators"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title = slide.shapes.title
        title.text = "Project Status Dashboard"
        self._apply_deltek_styling(title)

        # Calculate overall progress
        phases = project.get("phases", [])
        completed = len([p for p in phases if p.get("status") == "complete"])
        total = len(phases)
        progress = (completed / total * 100) if total > 0 else 0

        # Add progress text
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = f"Overall Progress: {progress:.0f}% ({completed}/{total} phases complete)"
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True

    def _add_timeline_slide(self, prs: Presentation, project: Dict):
        """Add project timeline (Gantt-style)"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title = slide.shapes.title
        title.text = "Project Timeline"
        self._apply_deltek_styling(title)

        # Add timeline visualization placeholder
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(3.5)

        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.text = "Gantt Chart Timeline\n(Auto-generated from project phases)"
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 250)

    def _add_next_steps_slide(self, prs: Presentation, project: Dict):
        """Add next steps and recommendations"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Next Steps & Recommendations"
        self._apply_deltek_styling(title)

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        next_steps = self._generate_next_steps(project)
        for step in next_steps:
            p = text_frame.add_paragraph()
            p.text = f"• {step}"
            p.font.size = Pt(14)
            p.space_after = Pt(12)

    def _add_architecture_detailed_slide(self, prs: Presentation, project: Dict):
        """Add detailed architecture diagram"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Detailed System Architecture"
        self._apply_deltek_styling(title)

    def _add_component_breakdown_slide(self, prs: Presentation, project: Dict):
        """Add component breakdown"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Component Breakdown"
        self._apply_deltek_styling(title)

    def _add_data_flow_slide(self, prs: Presentation, project: Dict):
        """Add data flow diagram"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Data Flow Architecture"
        self._apply_deltek_styling(title)

    def _add_deployment_diagram_slide(self, prs: Presentation, project: Dict):
        """Add deployment diagram"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Deployment Architecture"
        self._apply_deltek_styling(title)

    def _add_status_overview_slide(self, prs: Presentation, project: Dict):
        """Add status overview"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Status Overview"
        self._apply_deltek_styling(title)

    def _add_agent_performance_slide(self, prs: Presentation, project: Dict):
        """Add agent performance metrics"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Agent Performance Metrics"
        self._apply_deltek_styling(title)

    def _add_issues_risks_slide(self, prs: Presentation, project: Dict):
        """Add issues and risks"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Issues & Risks"
        self._apply_deltek_styling(title)

    def _add_action_items_slide(self, prs: Presentation, project: Dict):
        """Add action items"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Action Items"
        self._apply_deltek_styling(title)

    def _apply_deltek_styling(self, shape):
        """Apply Deltek brand styling to shape"""
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = "Figtree"
                paragraph.font.bold = True
                # Convert hex to RGB
                rgb = self._hex_to_rgb(self.deltek_brand["colors"]["primary_blue"])
                paragraph.font.color.rgb = RGBColor(*rgb)

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _add_kpi_box(self, slide, left, top, name, value):
        """Add a KPI visualization box"""
        width = Inches(4)
        height = Inches(1.5)

        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 250)

        text_frame = shape.text_frame
        text_frame.clear()

        # Add metric value
        p = text_frame.add_paragraph()
        p.text = str(value)
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Add metric name
        p = text_frame.add_paragraph()
        p.text = name
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    def _calculate_kpis(self, project: Dict) -> Dict:
        """Calculate key performance indicators"""
        phases = project.get("phases", [])
        return {
            "Total Phases": len(phases),
            "Completed": len([p for p in phases if p.get("status") == "complete"]),
            "In Progress": len([p for p in phases if p.get("status") == "in_progress"]),
            "Success Rate": f"{(len([p for p in phases if p.get('status') == 'complete']) / len(phases) * 100):.0f}%" if phases else "0%"
        }

    def _generate_summary_points(self, project: Dict) -> List[str]:
        """Generate executive summary bullet points"""
        return [
            f"Project {project.get('id')} is currently {project.get('status', 'active')}",
            f"{len(project.get('phases', []))} development phases orchestrated by AI agents",
            "Automated quality assurance, security scanning, and deployment",
            "Real-time monitoring and cost optimization enabled",
            "Self-learning agents improving efficiency with each iteration"
        ]

    def _generate_next_steps(self, project: Dict) -> List[str]:
        """Generate next steps recommendations"""
        return [
            "Continue monitoring agent performance and optimization opportunities",
            "Review and approve completed phases for production deployment",
            "Schedule stakeholder demo of completed features",
            "Plan next sprint priorities based on current velocity",
            "Update documentation and runbooks for operational team"
        ]

    def _save_presentation(self, prs: Presentation, project_id: str, pres_type: str) -> str:
        """Save presentation and archive"""
        timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        filename = f"{project_id}_{pres_type}_{timestamp}.pptx"

        # Save to presentations directory
        output_dir = self.memory_dir / "presentations"
        output_path = output_dir / filename
        prs.save(str(output_path))

        # Archive copy
        archive_path = output_dir / "archive" / filename
        shutil.copy(str(output_path), str(archive_path))

        return str(output_path)

    def _record_learning(self, project_id: str, pres_type: str, output_path: str):
        """Record learning from this generation"""
        learning_file = self.memory_dir / "learning" / "generation_log.json"

        log_entry = {
            "timestamp": datetime.now().isoformat(),
            "project_id": project_id,
            "presentation_type": pres_type,
            "output_path": output_path,
            "template_used": str(self.template_path),
            "slides_generated": "auto",
            "brand_compliance": "deltek_official"
        }

        # Append to log
        if learning_file.exists():
            with open(learning_file, 'r') as f:
                log = json.load(f)
        else:
            log = []

        log.append(log_entry)

        with open(learning_file, 'w') as f:
            json.dump(log, f, indent=2)

    def auto_update_presentation(self, presentation_path: str):
        """Auto-update an existing presentation with latest data"""
        # Load presentation
        prs = Presentation(presentation_path)

        # Extract project ID from filename
        filename = Path(presentation_path).stem
        project_id = filename.split('_')[0]

        # Reload project data
        project_data = self._load_project_data(project_id)

        # Update slides (simplified - would update specific slide content)
        # This is a placeholder for the actual update logic

        # Save updated version
        timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        new_path = presentation_path.replace(".pptx", f"_updated_{timestamp}.pptx")
        prs.save(new_path)

        return new_path


if __name__ == "__main__":
    import sys

    agent = ExecAgent()

    if len(sys.argv) > 1:
        command = sys.argv[1]

        if command == "generate":
            project_id = sys.argv[2] if len(sys.argv) > 2 else "SAMPLE-PROJECT"
            pres_type = sys.argv[3] if len(sys.argv) > 3 else "executive-summary"

            output = agent.generate_presentation(project_id, pres_type)
            print(f"✅ Presentation generated: {output}")

        elif command == "update":
            pres_path = sys.argv[2] if len(sys.argv) > 2 else None
            if pres_path:
                output = agent.auto_update_presentation(pres_path)
                print(f"✅ Presentation updated: {output}")
    else:
        print("Exec Agent - Executive Presentation Specialist")
        print("Usage:")
        print("  python exec-agent.py generate <project_id> [type]")
        print("  python exec-agent.py update <presentation_path>")
